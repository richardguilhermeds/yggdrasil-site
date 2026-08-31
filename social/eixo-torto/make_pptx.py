# -*- coding: utf-8 -*-
"""Monta os 6 posts da série "Eixo Torto" como PPTX de 1080x1350 (1 slide cada),
prontos para importar no Canva como designs editáveis.

Uso: python3 make_pptx.py <dir_fontes> <dir_charts> <dir_saida>
"""
import sys
from pathlib import Path

from PIL import ImageFont
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Pt

BG = "0E0F12"
INK = "F4F1EA"
MUTED = "9BA1AB"
ACCENT = "FF4D00"

W, H = 1080, 1350
MARGIN = 84
COL = W - 2 * MARGIN  # 912

FONT_DIR = Path(sys.argv[1])
CHARTS = Path(sys.argv[2])
OUT = Path(sys.argv[3])
OUT.mkdir(parents=True, exist_ok=True)

F_AB = FONT_DIR / "ArchivoBlack-Regular.ttf"
F_SG = FONT_DIR / "SpaceGrotesk-Regular.ttf"
F_SGB = FONT_DIR / "SpaceGrotesk-Bold.ttf"


def px(v):
    return Emu(int(round(v * 9525)))


def n_lines(text, ttf, size_px, max_w):
    """Número de linhas após word-wrap na largura max_w (px lógicos)."""
    font = ImageFont.truetype(str(ttf), size_px)
    lines, cur = 1, ""
    for word in text.split():
        trial = (cur + " " + word).strip()
        if font.getlength(trial) <= max_w or not cur:
            cur = trial
        else:
            lines += 1
            cur = word
    return lines


def add_text(slide, text, x, y, w, size_px, ttf_name, color, bold=False,
             line_spacing=1.0, align=PP_ALIGN.LEFT, height_px=None):
    lines = n_lines(text, {"Archivo Black": F_AB, "Space Grotesk": F_SGB if bold else F_SG}[ttf_name],
                    size_px, w)
    h = height_px or int(lines * size_px * line_spacing + size_px * 0.45)
    box = slide.shapes.add_textbox(px(x), px(y), px(w), px(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    for m in ("margin_left", "margin_right", "margin_top", "margin_bottom"):
        setattr(tf, m, 0)
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = line_spacing
    run = p.add_run()
    run.text = text
    f = run.font
    f.name = ttf_name
    f.size = Pt(size_px * 0.75)
    f.bold = bold
    f.color.rgb = RGBColor.from_string(color)
    return y + lines * int(size_px * line_spacing)


def add_image(slide, path, x, y, w):
    pic = slide.shapes.add_picture(str(path), px(x), px(y), width=px(w))
    return y + pic.height / 9525


def new_slide(prs, title):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    from pptx.enum.shapes import MSO_SHAPE
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, px(0), px(0), px(W), px(H))
    rect.fill.solid()
    rect.fill.fore_color.rgb = RGBColor.from_string(BG)
    rect.line.fill.background()
    rect.shadow.inherit = False
    return slide


def footer(slide):
    y = H - MARGIN - 44
    add_image(slide, CHARTS / "marca_eixo_torto.png", MARGIN, y + 2, 40)
    add_text(slide, "Richard Guilherme · ciência de dados & risco de crédito",
             MARGIN + 58, y + 8, COL - 58, 24, "Space Grotesk", MUTED)


POSTS = [
    dict(
        arquivo="00_abertura.pptx", titulo="Eixo Torto 00 — Abertura da série",
        kicker="NOVA SÉRIE",
        h1="Todo gráfico conta uma história. Alguns contam mentira.",
        h1_size=84,
        hero=True,
        body=("Eixo cortado, escala esperta, janela conveniente, acumulado "
              "eterno: os truques que transformam dado honesto em manchete torta."),
        cta="→ Um crime por semana. Salva e segue pra não cair no próximo.",
    ),
    dict(
        arquivo="01_eixo_cortado.pptx", titulo="Eixo Torto 01 — O eixo cortado",
        kicker="EIXO TORTO — CRIME 01",
        h1="O crescimento que só existe no eixo",
        h1_size=78, chart="crime01_eixo_cortado.png",
        caption=("Mesmos números, mesma variação: +0,9 ponto. "
                 "O que muda é onde o eixo começa."),
        cta="→ Salva pra próxima reunião de resultados.",
    ),
    dict(
        arquivo="02_dois_eixos.pptx", titulo="Eixo Torto 02 — Dois eixos",
        kicker="EIXO TORTO — CRIME 02",
        h1="Dois eixos provam qualquer tese",
        h1_size=78, chart="crime02_dois_eixos.png",
        caption=("Cada série com seu eixo, cada eixo com a escala escolhida a "
                 "dedo — e qualquer dupla de curvas “combina”."),
        cta="→ Comenta o pior eixo duplo que você já viu num dashboard.",
    ),
    dict(
        arquivo="03_raio_area.pptx", titulo="Eixo Torto 03 — Raio vs. área",
        kicker="EIXO TORTO — CRIME 03",
        h1="O dobro que parece o quádruplo",
        h1_size=78, chart="crime03_raio_area.png",
        caption=("O valor dobrou, mas dobrar o raio multiplica a área por "
                 "quatro — e o olho lê área."),
        cta="→ Manda com carinho pra quem ama infográfico de bolha.",
    ),
    dict(
        arquivo="04_janela.pptx", titulo="Eixo Torto 04 — A janela conveniente",
        kicker="EIXO TORTO — CRIME 04",
        h1="Escolha a janela, escolha a manchete",
        h1_size=78, chart="crime04_janela.png",
        caption=("A mesma série sustenta “+18% no trimestre” e “−40% em três "
                 "anos”. Depende só de onde você corta."),
        cta="→ Antes de acreditar, pergunta: desde quando?",
    ),
    dict(
        arquivo="05_acumulado.pptx", titulo="Eixo Torto 05 — O acumulado",
        kicker="EIXO TORTO — CRIME 05",
        h1="Acumulado nunca cai",
        h1_size=78, chart="crime05_acumulado.png",
        caption=("O total acumulado cresce até no mês em que quase ninguém "
                 "chegou. A curva de baixo é o produto."),
        cta="→ Salva a série inteira antes da próxima all-hands.",
    ),
]


def build(post):
    prs = Presentation()
    prs.slide_width = px(W)
    prs.slide_height = px(H)
    slide = new_slide(prs, post["titulo"])

    add_text(slide, post["kicker"], MARGIN, MARGIN, COL, 26, "Space Grotesk",
             ACCENT, bold=True)
    y = add_text(slide, post["h1"], MARGIN, MARGIN + 52, COL, post["h1_size"],
                 "Archivo Black", INK, line_spacing=1.04)
    y += 34

    if post.get("hero"):
        add_image(slide, CHARTS / "marca_eixo_torto_hero.png",
                  W - MARGIN - 330, y + 40, 330)
        y = add_text(slide, post["body"], MARGIN, y + 70, 540, 36,
                     "Space Grotesk", INK, line_spacing=1.25)
        y_cta = H - MARGIN - 44 - 150
        add_text(slide, post["cta"], MARGIN, y_cta, COL, 32, "Space Grotesk",
                 ACCENT, bold=True, line_spacing=1.2)
    else:
        y = add_image(slide, CHARTS / post["chart"], (W - 936) / 2, y + 6, 936)
        y = add_text(slide, post["caption"], MARGIN, y + 30, COL, 31,
                     "Space Grotesk", INK, line_spacing=1.25)
        add_text(slide, post["cta"], MARGIN, y + 26, COL, 29, "Space Grotesk",
                 ACCENT, bold=True, line_spacing=1.2)

    footer(slide)
    prs.save(OUT / post["arquivo"])
    print("ok", post["arquivo"])


if __name__ == "__main__":
    for p in POSTS:
        build(p)
