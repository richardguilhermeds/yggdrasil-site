# -*- coding: utf-8 -*-
"""Cinco layouts alternativos do verbete SAFRA (1080x1350), com o gráfico
desenhado em FORMAS NATIVAS (freeform/shapes) em vez de imagem — assim cada
curva, barra e marca fica selecionável e recolorível no Canva.

Gera .pptx (import no Canva) e .png (prévia local) do mesmo layout.

Uso: python3 make_alt_safra.py <dir_fontes> <dir_pptx> <dir_preview>
"""
import sys
from pathlib import Path

import numpy as np
from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Pt

W, H = 1080, 1350

# paleta escura (padrão da série)
BG = "0E0F12"
INK = "F4F1EA"
MUTED = "9BA1AB"
GRIDC = "22262C"
DIM = "2A2F36"
OLD = "5A6169"
ACC = "BEF264"
PANEL = "171B20"
# paleta clara (layout E, "caderno")
PAPER = "F4F1EA"
EINK = "17181C"
EMUT = "8A8578"
EGRID = "DDD6C8"
EDIM = "B9B3A6"
RED = "D9482B"

FONT_DIR = Path(sys.argv[1])
OUT_PPTX = Path(sys.argv[2])
OUT_PNG = Path(sys.argv[3])
OUT_PPTX.mkdir(parents=True, exist_ok=True)
OUT_PNG.mkdir(parents=True, exist_ok=True)

TTF = {
    ("Archivo Black", False): FONT_DIR / "ArchivoBlack-Regular.ttf",
    ("Space Grotesk", False): FONT_DIR / "SpaceGrotesk-Regular.ttf",
    ("Space Grotesk", True): FONT_DIR / "SpaceGrotesk-Bold.ttf",
    ("Caveat", False): FONT_DIR / "Caveat-600.ttf",
    ("Caveat", True): FONT_DIR / "Caveat-600.ttf",
}

# ---------------------------------------------------------------- dados
M_OLD = np.arange(0, 18.5, 0.5)
M_BAD = np.arange(0, 10.5, 0.5)
OLDS = [4.4 * (1 - np.exp(-0.22 * M_OLD)),
        3.8 * (1 - np.exp(-0.22 * M_OLD)),
        4.9 * (1 - np.exp(-0.22 * M_OLD))]
BAD = 8.6 * (1 - np.exp(-0.22 * M_BAD))  # 7,6% aos 10 meses

NOME = "análise vintage, a carteira contada por geração"
DEF = ("Cada mês de originação vira uma geração acompanhada separadamente: "
       "quanto da safra deu problema 3, 6, 12 meses após a concessão.")
PEG = ("Na média a carteira parece bem — até a safra ruim fazer aniversário. "
       "Média de carteira esconde geração problema.")
CTA = "→ Salva: é o primeiro gráfico quando a inadimplência “surpreende”."
KICKER = "GLOSSÁRIO DE RISCO — VERBETE 05"
BYLINE = "Richard Guilherme · ciência de dados & risco de crédito"


def hexc(h):
    return tuple(int(h[i:i + 2], 16) for i in (0, 2, 4))


def wrap(text, font, max_w):
    lines, cur = [], ""
    for word in text.split():
        trial = (cur + " " + word).strip()
        if font.getlength(trial) <= max_w or not cur:
            cur = trial
        else:
            lines.append(cur)
            cur = word
    lines.append(cur)
    return lines


def ribbon_pts(points, width):
    """Transforma uma polilinha em polígono-fita (para curva recolorível)."""
    p = np.asarray(points, dtype=float)
    d = np.gradient(p, axis=0)
    n = np.stack([-d[:, 1], d[:, 0]], axis=1)
    n /= np.linalg.norm(n, axis=1, keepdims=True)
    top = p + n * width / 2
    bot = p - n * width / 2
    return [tuple(q) for q in np.vstack([top, bot[::-1]])]


class Png:
    """Prévia local; desenha em 2x e reduz, para suavizar as formas."""
    S = 2

    def __init__(self, bg):
        s = self.S
        self.img = Image.new("RGBA", (W * s, H * s), hexc(bg) + (255,))
        self.d = ImageDraw.Draw(self.img)

    def rect(self, x, y, w, h, color, r=0):
        s = self.S
        box = (x * s, y * s, (x + w) * s, (y + h) * s)
        if r:
            self.d.rounded_rectangle(box, radius=r * s, fill=hexc(color))
        else:
            self.d.rectangle(box, fill=hexc(color))

    def oval(self, x, y, w, h, color):
        s = self.S
        self.d.ellipse((x * s, y * s, (x + w) * s, (y + h) * s), fill=hexc(color))

    def oval_dashed(self, x, y, w, h, color, weight):
        s = self.S
        box = (x * s, y * s, (x + w) * s, (y + h) * s)
        for a in range(0, 360, 24):
            self.d.arc(box, a, a + 14, fill=hexc(color), width=int(weight * s))

    def poly(self, points, color):
        s = self.S
        self.d.polygon([(px * s, py * s) for px, py in points], fill=hexc(color))

    def ribbon(self, points, width, color):
        self.poly(ribbon_pts(points, width), color)

    def text(self, t, x, y, w, size, fam, color, bold=False, ls=1.18,
             align="l"):
        s = self.S
        font = ImageFont.truetype(str(TTF[(fam, bold)]), size * s)
        for ln in wrap(t, font, w * s):
            lw = font.getlength(ln)
            xx = x * s + (w * s - lw) * {"l": 0, "c": 0.5, "r": 1.0}[align]
            self.d.text((xx, y * s), ln, font=font, fill=hexc(color))
            y += int(size * ls)
        return y

    def save(self, name):
        self.img.resize((W, H), Image.LANCZOS).convert("RGB").save(
            OUT_PNG / f"{name}.png")


class Pptx:
    def __init__(self, bg):
        self.prs = Presentation()
        self.prs.slide_width = Emu(W * 9525)
        self.prs.slide_height = Emu(H * 9525)
        self.slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self.rect(0, 0, W, H, bg)

    @staticmethod
    def _no_line(shape):
        shape.line.fill.background()
        shape.shadow.inherit = False

    def _fill(self, shape, color):
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor.from_string(color)
        self._no_line(shape)

    def rect(self, x, y, w, h, color, r=0):
        kind = MSO_SHAPE.ROUNDED_RECTANGLE if r else MSO_SHAPE.RECTANGLE
        shp = self.slide.shapes.add_shape(
            kind, Emu(int(x * 9525)), Emu(int(y * 9525)),
            Emu(int(w * 9525)), Emu(int(h * 9525)))
        if r:
            shp.adjustments[0] = min(0.5, r / min(w, h))
        self._fill(shp, color)

    def oval(self, x, y, w, h, color):
        shp = self.slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Emu(int(x * 9525)), Emu(int(y * 9525)),
            Emu(int(w * 9525)), Emu(int(h * 9525)))
        self._fill(shp, color)

    def oval_dashed(self, x, y, w, h, color, weight):
        shp = self.slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Emu(int(x * 9525)), Emu(int(y * 9525)),
            Emu(int(w * 9525)), Emu(int(h * 9525)))
        shp.fill.background()
        shp.shadow.inherit = False
        shp.line.color.rgb = RGBColor.from_string(color)
        shp.line.width = Pt(weight * 0.75)
        shp.line.dash_style = MSO_LINE_DASH_STYLE.DASH

    def poly(self, points, color):
        b = self.slide.shapes.build_freeform(points[0][0], points[0][1],
                                             scale=9525)
        b.add_line_segments(points[1:], close=True)
        self._fill(b.convert_to_shape(), color)

    def ribbon(self, points, width, color):
        self.poly(ribbon_pts(points, width), color)

    def text(self, t, x, y, w, size, fam, color, bold=False, ls=1.18,
             align="l"):
        font = ImageFont.truetype(str(TTF[(fam, bold)]), size)
        lines = wrap(t, font, w)
        h = int(len(lines) * size * ls + size * 0.5)
        box = self.slide.shapes.add_textbox(
            Emu(int(x * 9525)), Emu(int(y * 9525)),
            Emu(int(w * 9525)), Emu(int(h * 9525)))
        tf = box.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.TOP
        for m in ("margin_left", "margin_right", "margin_top", "margin_bottom"):
            setattr(tf, m, 0)
        p = tf.paragraphs[0]
        p.alignment = {"l": PP_ALIGN.LEFT, "c": PP_ALIGN.CENTER,
                       "r": PP_ALIGN.RIGHT}[align]
        p.line_spacing = ls
        run = p.add_run()
        run.text = t
        run.font.name = fam
        run.font.size = Pt(size * 0.75)
        run.font.bold = bold
        run.font.color.rgb = RGBColor.from_string(color)
        return y + len(lines) * int(size * ls)

    def save(self, name):
        self.prs.save(OUT_PPTX / f"{name}.pptx")


def to_px(m, v, box, mmax=18.5, vmax=9.0):
    """Mapeia (mês, %) para pixels dentro de box=(x, y, w, h)."""
    x, y, w, h = box
    return x + m / mmax * w, y + h - v / vmax * h


def curvas(r, box, w_old=6, w_bad=10, old_color=OLD, bad_color=ACC,
           mmax=18.5, vmax=9.0):
    for serie in OLDS:
        pts = [to_px(m, v, box, mmax, vmax) for m, v in zip(M_OLD, serie)]
        r.ribbon(pts, w_old, old_color)
    pts = [to_px(m, v, box, mmax, vmax) for m, v in zip(M_BAD, BAD)]
    r.ribbon(pts, w_bad, bad_color)
    return to_px(M_BAD[-1], BAD[-1], box, mmax, vmax)


def marca(r, x, y, color, s=1.0):
    r.oval(x, y, 11 * s, 11 * s, color)
    r.oval(x, y + 20 * s, 11 * s, 11 * s, color)
    r.rect(x + 20 * s, y, 34 * s, 11 * s, color, r=5 * s)
    r.rect(x + 20 * s, y + 20 * s, 34 * s, 11 * s, color, r=5 * s)


def rodape(r, color_marca, color_txt, x=84, y=H - 84 - 40, texto=BYLINE):
    marca(r, x, y + 6, color_marca)
    r.text(texto, x + 74, y + 8, W - x - 74 - 84, 24, "Space Grotesk",
           color_txt)


def aviso(nome, bottom, limite):
    if bottom > limite:
        print(f"  AVISO {nome}: conteúdo em y={bottom:.0f} passa de {limite}")


# ================================================================ A
def layout_a(r):
    # gráfico full-bleed no topo, texto embaixo
    box = (-30, 120, 1140, 540)
    for gy in (200, 360, 520):
        r.rect(-2, gy, W + 4, 2, GRIDC)
    end = curvas(r, box)
    r.rect(-2, 660, W + 4, 2, "3A3F47")
    r.text(KICKER, 84, 84, 912, 26, "Space Grotesk", ACC, bold=True)
    r.text("7,6% aos 10 meses", end[0] + 22, end[1] - 40, 380, 26,
           "Space Grotesk", ACC, bold=True)
    r.text("a safra nova", end[0] + 22, end[1] - 6, 380, 22,
           "Space Grotesk", MUTED)
    r.text("safras anteriores", 848, 344, 200, 22, "Space Grotesk", MUTED)
    r.text("originação", 84, 676, 300, 22, "Space Grotesk", MUTED)
    r.text("18 meses", W - 84 - 300, 676, 300, 22, "Space Grotesk", MUTED,
           align="r")

    y = r.text("SAFRA", 84, 740, 912, 118, "Archivo Black", INK, ls=1.0)
    y = r.text(NOME, 84, y + 10, 912, 28, "Space Grotesk", MUTED)
    y = r.text(DEF, 84, y + 18, 912, 29, "Space Grotesk", INK, ls=1.28)
    y = r.text(PEG, 84, y + 18, 912, 29, "Space Grotesk", ACC, bold=True,
               ls=1.25)
    y = r.text(CTA, 84, y + 16, 912, 24, "Space Grotesk", MUTED, bold=True)
    rodape(r, ACC, MUTED)
    aviso("A", y, H - 84 - 44 - 12)


# ================================================================ B
def layout_b(r):
    # cartão de dicionário sobre campo de cor
    r.rect(56, 56, W - 112, H - 112, BG, r=40)
    x, w = 128, 824
    r.text(KICKER, x, 128, w, 24, "Space Grotesk", ACC, bold=True)
    r.text("sa·fra", x, 174, w, 116, "Archivo Black", INK, ls=1.0)
    r.text("substantivo feminino · análise vintage · risco de crédito",
           x, 318, w, 27, "Space Grotesk", MUTED)
    r.rect(x, 372, w, 3, DIM)

    y = r.text("1. " + DEF, x, 404, w, 29, "Space Grotesk", INK, ls=1.3)

    box = (x, 556, w, 320)
    for gy in (556 + 90, 556 + 200):
        r.rect(x, gy, w, 2, GRIDC)
    end = curvas(r, box, w_old=5, w_bad=9)
    r.rect(x, 876, w, 2, "3A3F47")
    r.text("7,6%", end[0] + 16, end[1] - 14, 120, 26, "Space Grotesk", ACC,
           bold=True)
    r.text("3,9%", x + w - 96, 750, 96, 24, "Space Grotesk", MUTED, bold=True)

    y = r.text("2. Motivo pelo qual a média da carteira engana: no agregado "
               "tudo parece bem — até a safra ruim fazer aniversário.",
               x, 926, w, 29, "Space Grotesk", ACC, bold=True, ls=1.28)
    y = r.text("→ Salva o verbete. Sai um por semana.", x, y + 22, w, 24,
               "Space Grotesk", MUTED, bold=True)
    rodape(r, ACC, MUTED, x=x, y=H - 128 - 40,
           texto="Richard Guilherme · risco de crédito")
    aviso("B", y, H - 128 - 40 - 12)


# ================================================================ C
def layout_c(r):
    # split vertical: tipografia à esquerda, painel de dado à direita
    r.rect(470, 0, W - 470, H, PANEL)
    lx, lw = 84, 340
    r.text(KICKER, lx, 84, lw, 23, "Space Grotesk", ACC, bold=True)
    y = r.text("SA FRA", lx, 182, 300, 148, "Archivo Black", INK, ls=0.98)
    y = r.text(NOME, lx, y + 14, lw, 26, "Space Grotesk", MUTED, ls=1.25)
    y = r.text(DEF, lx, y + 24, lw, 26, "Space Grotesk", INK, ls=1.3)
    y = r.text(PEG, lx, y + 22, lw, 26, "Space Grotesk", ACC, bold=True,
               ls=1.28)
    y = r.text("→ Salva pra depois.", lx, y + 20, lw, 23, "Space Grotesk",
               MUTED, bold=True)
    marca(r, lx, H - 84 - 34, ACC)
    r.text("Richard Guilherme", lx + 74, H - 84 - 26, lw - 74, 23,
           "Space Grotesk", MUTED)

    px_, pw = 526, 498
    r.text("7,6%", px_, 84, pw, 96, "Archivo Black", ACC, ls=1.0)
    r.text("da safra nova já com problema aos 10 meses — o triplo das "
           "anteriores", px_, 196, pw, 24, "Space Grotesk", MUTED, ls=1.3)

    box = (px_, 320, pw, 750)
    for gy in (320 + 190, 320 + 380, 320 + 570):
        r.rect(px_, gy, pw, 2, "20242A")
    curvas(r, box, w_old=6, w_bad=10)
    r.rect(px_, 1070, pw, 2, "3A3F47")
    r.text("originação", px_, 1086, 240, 21, "Space Grotesk", MUTED)
    r.text("18 meses", px_ + pw - 240, 1086, 240, 21, "Space Grotesk", MUTED,
           align="r")
    r.oval(px_, 1150, 14, 14, ACC)
    r.text("a safra nova", px_ + 26, 1146, 200, 22, "Space Grotesk", INK)
    r.oval(px_ + 240, 1150, 14, 14, OLD)
    r.text("safras anteriores", px_ + 266, 1146, 232, 22, "Space Grotesk",
           MUTED)
    aviso("C", y, H - 84 - 34 - 12)


# ================================================================ D
def layout_d(r):
    # números primeiro: o contraste vira o herói
    r.text(KICKER, 84, 84, 912, 26, "Space Grotesk", ACC, bold=True)

    r.text("3,9%", 84, 170, 400, 150, "Archivo Black", MUTED, ls=1.0)
    r.text("safras anteriores,\nproblema aos 10 meses", 84, 348, 380, 26,
           "Space Grotesk", MUTED, ls=1.3)
    box = (84, 470, 380, 120)
    serie = 4.37 * (1 - np.exp(-0.22 * M_BAD))
    pts = [to_px(m, v, box, mmax=10.5, vmax=9.0) for m, v in zip(M_BAD, serie)]
    r.ribbon(pts, 7, OLD)

    r.text("vs", 512, 240, 56, 30, "Space Grotesk", MUTED, align="c")

    r.text("7,6%", 596, 170, 400, 150, "Archivo Black", ACC, ls=1.0)
    r.text("a safra nova, aos mesmos 10 meses", 596, 348, 380, 26,
           "Space Grotesk", INK, ls=1.3)
    box = (596, 470, 380, 120)
    pts = [to_px(m, v, box, mmax=10.5, vmax=9.0) for m, v in zip(M_BAD, BAD)]
    r.ribbon(pts, 9, ACC)

    r.rect(84, 648, 912, 2, GRIDC)

    y = r.text("SAFRA", 84, 692, 912, 96, "Archivo Black", INK, ls=1.0)
    y = r.text(NOME, 84, y + 8, 912, 28, "Space Grotesk", MUTED)
    y = r.text(DEF, 84, y + 18, 912, 29, "Space Grotesk", INK, ls=1.28)
    y = r.text(PEG, 84, y + 18, 912, 29, "Space Grotesk", ACC, bold=True,
               ls=1.25)
    y = r.text(CTA, 84, y + 16, 912, 24, "Space Grotesk", MUTED, bold=True)
    rodape(r, ACC, MUTED)
    aviso("D", y, H - 84 - 44 - 12)


# ================================================================ E
def layout_e(r):
    # caderno do analista: fundo claro, anotação à mão
    r.text(KICKER, 84, 84, 912, 24, "Space Grotesk", RED, bold=True)
    r.text("SAFRA", 84, 128, 912, 120, "Archivo Black", EINK, ls=1.0)
    r.text(NOME, 84, 276, 912, 28, "Space Grotesk", EMUT)

    box = (84, 400, 912, 500)
    for gy in (400 + 125, 400 + 250, 400 + 375):
        r.rect(84, gy, 912, 2, EGRID)
    end = curvas(r, box, w_old=6, w_bad=10, old_color=EDIM, bad_color=EINK)
    r.rect(84, 900, 912, 2, EINK)
    r.text("originação", 84, 918, 300, 21, "Space Grotesk", EMUT)
    r.text("18 meses", W - 84 - 300, 918, 300, 21, "Space Grotesk", EMUT,
           align="r")

    r.oval_dashed(end[0] - 78, end[1] - 52, 156, 104, RED, 4)
    r.text("ainda jovem, e já pior que todo mundo", 560, 300, 420, 42,
           "Caveat", RED, ls=1.05)
    r.ribbon([(700, 392), (660, 410), (622, 432)], 4, RED)
    r.poly([(628, 442), (612, 426), (634, 421)], RED)

    y = r.text(DEF, 84, 964, 912, 29, "Space Grotesk", EINK, ls=1.28)
    y = r.text("na média parece tudo bem. a média mente.", 84, y + 16, 912,
               46, "Caveat", RED)
    y = r.text(CTA, 84, y + 12, 912, 24, "Space Grotesk", EINK, bold=True)
    rodape(r, RED, EMUT)
    aviso("E", y, H - 84 - 44 - 12)


LAYOUTS = [
    ("alt_a_fullbleed", BG, layout_a),
    ("alt_b_cartao", ACC, layout_b),
    ("alt_c_split", BG, layout_c),
    ("alt_d_numero", BG, layout_d),
    ("alt_e_caderno", PAPER, layout_e),
]

if __name__ == "__main__":
    for nome, bg, fn in LAYOUTS:
        for r in (Pptx(bg), Png(bg)):
            fn(r)
            r.save(nome)
        print("ok", nome)
