# -*- coding: utf-8 -*-
"""Monta os 6 posts da série "Glossário de Risco" (1080x1350, 1 slide cada)
em PPTX (para importar no Canva) e PNG (prévia local), a partir do mesmo layout.

Uso: python3 make_posts.py <dir_fontes> <dir_charts> <dir_pptx> <dir_preview>
"""
import sys
from pathlib import Path

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Pt

BG = "0E0F12"
INK = "F4F1EA"
MUTED = "9BA1AB"
ACCENT = "BEF264"  # verde-lima da série

W, H = 1080, 1350
MARGIN = 84
COL = W - 2 * MARGIN
FOOTER_Y = H - MARGIN - 44

FONT_DIR = Path(sys.argv[1])
CHARTS = Path(sys.argv[2])
OUT_PPTX = Path(sys.argv[3])
OUT_PNG = Path(sys.argv[4])
OUT_PPTX.mkdir(parents=True, exist_ok=True)
OUT_PNG.mkdir(parents=True, exist_ok=True)

TTF = {
    ("Archivo Black", False): FONT_DIR / "ArchivoBlack-Regular.ttf",
    ("Space Grotesk", False): FONT_DIR / "SpaceGrotesk-Regular.ttf",
    ("Space Grotesk", True): FONT_DIR / "SpaceGrotesk-Bold.ttf",
}


def px(v):
    return Emu(int(round(v * 9525)))


def hexc(h):
    return tuple(int(h[i:i + 2], 16) for i in (0, 2, 4))


def wrap(text, font, max_w):
    lines, cur = [], ""
    for word in text.split():
        trial = (cur + " " + word).strip()
        if font.getlength(trial) <= max_w or not cur:
            cur = trial
        else:
            lines.append(cur)
            cur = word
    lines.append(cur)
    return lines


POSTS = [
    dict(
        arquivo="00_abertura", titulo="Glossário de Risco 00 — Abertura",
        kicker="NOVA SÉRIE",
        hero=True,
        h1="O glossário visual do risco de crédito.",
        body=("Um termo por post, um diagrama por termo, zero juridiquês. "
              "PD, LGD, KS, PSI, safra — e o que mais aparecer no comitê "
              "com cara séria."),
        cta="→ Um verbete por semana. Salva a coleção.",
    ),
    dict(
        arquivo="01_pd", titulo="Glossário de Risco 01 — PD",
        kicker="GLOSSÁRIO DE RISCO — VERBETE 01",
        termo="PD", nome="probabilidade de default",
        chart="verbete_pd.png",
        definicao=("A chance de o cliente entrar em default dentro de um "
                   "horizonte definido: 12 meses — ou a vida toda do contrato "
                   "(lifetime, no mundo IFRS 9/4.966)."),
        pegadinha=("PD de 2% não é “quase nunca”. É 1 em 50 — e o preço do "
                   "crédito inteiro se apoia nesse número."),
        cta="→ Salva pra próxima vez que disserem “esse cliente não tem risco”.",
    ),
    dict(
        arquivo="02_lgd", titulo="Glossário de Risco 02 — LGD",
        kicker="GLOSSÁRIO DE RISCO — VERBETE 02",
        termo="LGD", nome="perda dado o default",
        chart="verbete_lgd.png",
        definicao=("De tudo que estava exposto no momento do default, a fração "
                   "que não volta — líquida de recuperações, garantias e do "
                   "custo de recuperar."),
        pegadinha=("Default não é perder tudo. Mas o que volta, volta tarde — "
                   "e dinheiro atrasado vale menos: o desconto faz parte da conta."),
        cta="→ Manda pra quem acha que garantia resolve 100%.",
    ),
    dict(
        arquivo="03_ks", titulo="Glossário de Risco 03 — KS",
        kicker="GLOSSÁRIO DE RISCO — VERBETE 03",
        termo="KS", nome="Kolmogorov–Smirnov, a separação do modelo",
        chart="verbete_ks.png",
        definicao=("A maior distância entre as curvas acumuladas de bons e "
                   "maus ao longo do score. Quanto maior, melhor o modelo "
                   "separa os dois grupos."),
        pegadinha=("KS é medido no melhor ponto da régua — não garante a régua "
                   "inteira, nem que ela continua boa na safra que vem."),
        cta="→ Salva pro próximo comitê de modelos.",
    ),
    dict(
        arquivo="04_psi", titulo="Glossário de Risco 04 — PSI",
        kicker="GLOSSÁRIO DE RISCO — VERBETE 04",
        termo="PSI", nome="population stability index, o termômetro de drift",
        chart="verbete_psi.png",
        definicao=("Compara a distribuição do score de hoje com a da época do "
                   "desenvolvimento: mede o quanto a população escorregou por "
                   "baixo do modelo."),
        pegadinha=("Modelo não quebra com estrondo. A população muda em "
                   "silêncio — e o score vira outro sem avisar."),
        cta="→ Comenta: de quanto em quanto tempo você olha o seu?",
    ),
    dict(
        arquivo="05_safra", titulo="Glossário de Risco 05 — Safra",
        kicker="GLOSSÁRIO DE RISCO — VERBETE 05",
        termo="SAFRA", nome="análise vintage, a carteira contada por geração",
        chart="verbete_safra.png",
        definicao=("Cada mês de originação vira uma geração acompanhada "
                   "separadamente: quanto da safra deu problema 3, 6, 12 meses "
                   "após a concessão."),
        pegadinha=("Na média a carteira parece bem — até a safra ruim fazer "
                   "aniversário. Média de carteira esconde geração problema."),
        cta="→ Salva: é o primeiro gráfico quando a inadimplência “surpreende”.",
    ),
]


class Png:
    def __init__(self):
        self.img = Image.new("RGBA", (W, H), hexc(BG) + (255,))
        self.d = ImageDraw.Draw(self.img)

    def text(self, t, x, y, w, size, fam, color, bold=False, ls=1.0):
        font = ImageFont.truetype(str(TTF[(fam, bold)]), size)
        lines = wrap(t, font, w)
        for ln in lines:
            self.d.text((x, y), ln, font=font, fill=hexc(color))
            y += int(size * ls)
        return y

    def image(self, path, x, y, w):
        pic = Image.open(path).convert("RGBA")
        pic = pic.resize((w, int(pic.height * w / pic.width)), Image.LANCZOS)
        self.img.alpha_composite(pic, (int(x), int(y)))
        return y + pic.height

    def save(self, name):
        self.img.convert("RGB").save(OUT_PNG / f"{name}.png")


class Pptx:
    def __init__(self):
        self.prs = Presentation()
        self.prs.slide_width = px(W)
        self.prs.slide_height = px(H)
        self.slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        rect = self.slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                           px(0), px(0), px(W), px(H))
        rect.fill.solid()
        rect.fill.fore_color.rgb = RGBColor.from_string(BG)
        rect.line.fill.background()
        rect.shadow.inherit = False

    def text(self, t, x, y, w, size, fam, color, bold=False, ls=1.0):
        font = ImageFont.truetype(str(TTF[(fam, bold)]), size)
        lines = wrap(t, font, w)
        h = int(len(lines) * size * ls + size * 0.45)
        box = self.slide.shapes.add_textbox(px(x), px(y), px(w), px(h))
        tf = box.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.TOP
        for m in ("margin_left", "margin_right", "margin_top", "margin_bottom"):
            setattr(tf, m, 0)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = ls
        run = p.add_run()
        run.text = t
        run.font.name = fam
        run.font.size = Pt(size * 0.75)
        run.font.bold = bold
        run.font.color.rgb = RGBColor.from_string(color)
        return y + len(lines) * int(size * ls)

    def image(self, path, x, y, w):
        pic = self.slide.shapes.add_picture(str(path), px(x), px(y), width=px(w))
        return y + pic.height / 9525

    def save(self, name):
        self.prs.save(OUT_PPTX / f"{name}.pptx")


def layout(r, post):
    r.text(post["kicker"], MARGIN, MARGIN, COL, 26, "Space Grotesk",
           ACCENT, bold=True)
    if post.get("hero"):
        y = r.text(post["h1"], MARGIN, MARGIN + 52, COL, 84,
                   "Archivo Black", INK, ls=1.04)
        r.image(CHARTS / "marca_glossario_hero.png", W - MARGIN - 400, y + 150, 400)
        r.text(post["body"], MARGIN, y + 90, 500, 36, "Space Grotesk", INK,
               ls=1.25)
        r.text(post["cta"], MARGIN, FOOTER_Y - 150, COL, 32, "Space Grotesk",
               ACCENT, bold=True, ls=1.2)
        bottom = FOOTER_Y - 60
    else:
        y = r.text(post["termo"], MARGIN, MARGIN + 44, COL, 140,
                   "Archivo Black", INK, ls=1.0)
        y = r.text(post["nome"], MARGIN, y + 16, COL, 30, "Space Grotesk",
                   MUTED, ls=1.2)
        y = r.image(CHARTS / post["chart"], (W - 936) / 2, y + 26, 936)
        y = r.text(post["definicao"], MARGIN, y + 28, COL, 30, "Space Grotesk",
                   INK, ls=1.28)
        y = r.text(post["pegadinha"], MARGIN, y + 22, COL, 30, "Space Grotesk",
                   ACCENT, bold=True, ls=1.25)
        bottom = r.text(post["cta"], MARGIN, y + 20, COL, 26, "Space Grotesk",
                        MUTED, bold=True, ls=1.2)
    fy = FOOTER_Y
    r.image(CHARTS / "marca_glossario.png", MARGIN, fy + 2, 40)
    r.text("Richard Guilherme · ciência de dados & risco de crédito",
           MARGIN + 58, fy + 8, COL - 58, 24, "Space Grotesk", MUTED)
    if bottom > fy - 16:
        print(f"  AVISO: conteúdo de {post['arquivo']} chega a y={bottom:.0f} "
              f"(rodapé em {fy})")


if __name__ == "__main__":
    for post in POSTS:
        for renderer in (Pptx(), Png()):
            layout(renderer, post)
            renderer.save(post["arquivo"])
        print("ok", post["arquivo"])
